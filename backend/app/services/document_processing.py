"""Extract text from PDF/DOCX/PPTX/TXT and split into overlapping chunks."""
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from app.services.storage import get_file


def extract_text(storage_path: str, file_type: str) -> str:
    file_bytes = get_file(storage_path)

    if file_type == "pdf":
        return _extract_pdf(file_bytes)
    if file_type == "docx":
        return _extract_docx(file_bytes)
    if file_type == "pptx":
        return _extract_pptx(file_bytes)
    if file_type == "txt":
        return file_bytes.decode("utf-8", errors="ignore")

    raise ValueError(f"Unsupported file_type: {file_type}")


def _extract_pdf(file_bytes: bytes) -> str:
    import io

    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(file_bytes: bytes) -> str:
    import io

    doc = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(file_bytes: bytes) -> str:
    import io

    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    """Simple word-based chunking with overlap. chunk_size/overlap are in words."""
    words = text.split()
    if not words:
        return []

    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        if end >= len(words):
            break
        start = end - overlap

    return chunks
